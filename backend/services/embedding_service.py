import io
import json
import os

import boto3

VECTOR_BUCKET_NAME = os.getenv("VECTOR_BUCKET_NAME", "sylli-vectors")
VECTOR_INDEX_NAME = os.getenv("VECTOR_INDEX_NAME", "materials-index")
# Use AWS_REGION_NAME (not AWS_REGION) — Lambda reserves AWS_REGION as a system env var
# that cannot be overridden in template.yaml; AWS_REGION_NAME is a custom env var.
REGION = os.getenv("AWS_REGION_NAME", "us-east-1")

bedrock = boto3.client("bedrock-runtime", region_name=REGION)
# s3vectors client initialized lazily to defer import-time error if boto3 version is wrong
_s3v = None


def _get_s3v():
    global _s3v
    if _s3v is None:
        _s3v = boto3.client("s3vectors", region_name=REGION)
    return _s3v


CHUNK_SIZE = 1500
CHUNK_OVERLAP = 200


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract plain text from PDF or PPTX bytes."""
    if file_type == "pdf":
        import fitz  # PyMuPDF — imported locally to isolate heavy dependency
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    elif file_type == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    else:
        raise ValueError(f"Unsupported file_type: {file_type!r}")


def chunk_text(text: str) -> list[str]:
    """Split text into overlapping fixed-size chunks. Returns empty list for empty/whitespace input."""
    if not text.strip():
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunks.append(text[start:end])
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return [c for c in chunks if c.strip()]


def embed_text(text: str) -> list[float]:
    """Embed text using Titan Text Embeddings V2. Returns list of 1024 floats."""
    response = bedrock.invoke_model(
        modelId="amazon.titan-embed-text-v2:0",
        body=json.dumps({"inputText": text, "dimensions": 1024}),
    )
    result = json.loads(response["body"].read())
    return result["embedding"]


def write_vectors_to_s3(
    material_id: str,
    user_id: str,
    week_number: int,
    chunks: list[str],
    embeddings: list[list[float]],
):
    """Write all chunk vectors to S3 Vectors with metadata for Phase 4 retrieval."""
    if not chunks:
        return
    s3v = _get_s3v()
    vectors = [
        {
            "key": f"{material_id}#chunk#{i}",
            "data": {"float32": embeddings[i]},
            "metadata": {
                "user_id": user_id,
                "material_id": material_id,
                "week_number": week_number,
                "chunk_index": i,
                "source_text": chunks[i][:500],  # truncated — S3 Vectors metadata size limit
            },
        }
        for i in range(len(chunks))
    ]
    s3v.put_vectors(
        vectorBucketName=VECTOR_BUCKET_NAME,
        indexName=VECTOR_INDEX_NAME,
        vectors=vectors,
    )
