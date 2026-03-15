"""
Tests for backend/services/embedding_service.py

Covers:
- chunk_text: empty input, chunk size constraints
- extract_text: unknown file_type raises ValueError
- embed_text: calls bedrock with correct modelId (mocked)
- write_vectors_to_s3: calls put_vectors with correct key format (mocked)
"""
import json
from unittest.mock import MagicMock, patch

import pytest


# ---------------------------------------------------------------------------
# chunk_text tests
# ---------------------------------------------------------------------------

def test_chunk_text_empty_string_returns_empty_list():
    from services.embedding_service import chunk_text
    assert chunk_text("") == []


def test_chunk_text_whitespace_only_returns_empty_list():
    from services.embedding_service import chunk_text
    assert chunk_text("   \n\t  ") == []


def test_chunk_text_short_text_returns_single_chunk():
    from services.embedding_service import chunk_text
    text = "Hello world"
    result = chunk_text(text)
    assert len(result) == 1
    assert result[0] == text


def test_chunk_text_all_chunks_within_size_limit():
    from services.embedding_service import chunk_text, CHUNK_SIZE
    # Create text larger than one chunk
    text = "A" * 5000
    result = chunk_text(text)
    assert len(result) > 1
    for chunk in result:
        assert len(chunk) <= CHUNK_SIZE


def test_chunk_text_no_whitespace_only_chunks():
    from services.embedding_service import chunk_text
    text = "word " * 1000
    result = chunk_text(text)
    for chunk in result:
        assert chunk.strip() != ""


# ---------------------------------------------------------------------------
# extract_text tests
# ---------------------------------------------------------------------------

def test_extract_text_raises_for_unknown_file_type():
    from services.embedding_service import extract_text
    with pytest.raises(ValueError, match="Unsupported file_type"):
        extract_text(b"some bytes", "docx")


def test_extract_text_raises_for_empty_file_type():
    from services.embedding_service import extract_text
    with pytest.raises(ValueError, match="Unsupported file_type"):
        extract_text(b"some bytes", "")


def test_extract_text_pdf_returns_string(tmp_path):
    """Create a minimal PDF in-memory and verify extract_text returns a string."""
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 72), "Hello PDF text")
    pdf_bytes = doc.tobytes()
    doc.close()

    from services.embedding_service import extract_text
    result = extract_text(pdf_bytes, "pdf")
    assert isinstance(result, str)
    assert "Hello PDF text" in result


def test_extract_text_pptx_returns_string():
    """Create a minimal PPTX in-memory and verify extract_text returns a string."""
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Hello PPTX text"
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    from services.embedding_service import extract_text
    result = extract_text(pptx_bytes, "pptx")
    assert isinstance(result, str)
    assert "Hello PPTX text" in result


# ---------------------------------------------------------------------------
# embed_text tests
# ---------------------------------------------------------------------------

def test_embed_text_calls_bedrock_with_correct_model():
    """embed_text must call bedrock.invoke_model with the Titan Embed V2 modelId."""
    fake_embedding = [0.1] * 1024
    fake_body = MagicMock()
    fake_body.read.return_value = json.dumps({"embedding": fake_embedding}).encode()
    fake_response = {"body": fake_body}

    with patch("services.embedding_service.bedrock") as mock_bedrock:
        mock_bedrock.invoke_model.return_value = fake_response
        from services.embedding_service import embed_text
        result = embed_text("some text")

    mock_bedrock.invoke_model.assert_called_once()
    call_kwargs = mock_bedrock.invoke_model.call_args
    assert call_kwargs.kwargs["modelId"] == "amazon.titan-embed-text-v2:0"
    body = json.loads(call_kwargs.kwargs["body"])
    assert body["inputText"] == "some text"
    assert body["dimensions"] == 1024
    assert result == fake_embedding


def test_embed_text_returns_list_of_floats():
    """embed_text returns the embedding list from the response."""
    fake_embedding = [float(i) / 1024 for i in range(1024)]
    fake_body = MagicMock()
    fake_body.read.return_value = json.dumps({"embedding": fake_embedding}).encode()

    with patch("services.embedding_service.bedrock") as mock_bedrock:
        mock_bedrock.invoke_model.return_value = {"body": fake_body}
        from services.embedding_service import embed_text
        result = embed_text("test")

    assert len(result) == 1024
    assert all(isinstance(v, float) for v in result)


# ---------------------------------------------------------------------------
# write_vectors_to_s3 tests
# ---------------------------------------------------------------------------

def test_write_vectors_to_s3_calls_put_vectors_with_correct_key_format():
    """write_vectors_to_s3 must call put_vectors with keys in pattern '{material_id}#chunk#{i}'."""
    mock_s3v = MagicMock()
    chunks = ["chunk zero text", "chunk one text"]
    embeddings = [[0.1] * 1024, [0.2] * 1024]

    with patch("services.embedding_service._get_s3v", return_value=mock_s3v):
        from services.embedding_service import write_vectors_to_s3
        write_vectors_to_s3("mat-001", "user-123", 3, chunks, embeddings)

    mock_s3v.put_vectors.assert_called_once()
    call_kwargs = mock_s3v.put_vectors.call_args.kwargs
    vectors = call_kwargs["vectors"]
    assert len(vectors) == 2
    assert vectors[0]["key"] == "mat-001#chunk#0"
    assert vectors[1]["key"] == "mat-001#chunk#1"


def test_write_vectors_to_s3_metadata_fields():
    """Each vector's metadata must include all required fields with correct values."""
    mock_s3v = MagicMock()
    chunks = ["a" * 600]  # longer than 500 to test truncation
    embeddings = [[0.5] * 1024]

    with patch("services.embedding_service._get_s3v", return_value=mock_s3v):
        from services.embedding_service import write_vectors_to_s3
        write_vectors_to_s3("mat-abc", "user-xyz", 5, chunks, embeddings)

    vectors = mock_s3v.put_vectors.call_args.kwargs["vectors"]
    meta = vectors[0]["metadata"]
    assert meta["user_id"] == "user-xyz"
    assert meta["material_id"] == "mat-abc"
    assert meta["week_number"] == 5
    assert meta["chunk_index"] == 0
    assert len(meta["source_text"]) <= 500


def test_write_vectors_to_s3_no_call_when_empty_chunks():
    """write_vectors_to_s3 must NOT call put_vectors if there are no chunks."""
    mock_s3v = MagicMock()

    with patch("services.embedding_service._get_s3v", return_value=mock_s3v):
        from services.embedding_service import write_vectors_to_s3
        write_vectors_to_s3("mat-000", "user-000", 1, [], [])

    mock_s3v.put_vectors.assert_not_called()


def test_write_vectors_to_s3_vector_data_format():
    """Each vector's data must be wrapped as {'float32': embedding}."""
    mock_s3v = MagicMock()
    embedding = [0.3] * 1024
    chunks = ["test chunk"]
    embeddings = [embedding]

    with patch("services.embedding_service._get_s3v", return_value=mock_s3v):
        from services.embedding_service import write_vectors_to_s3
        write_vectors_to_s3("mat-xyz", "user-aaa", 2, chunks, embeddings)

    vectors = mock_s3v.put_vectors.call_args.kwargs["vectors"]
    assert vectors[0]["data"] == {"float32": embedding}
